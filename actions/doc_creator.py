"""
doc_creator.py — JARVIS Universal Document & Visual Creator

Creates from a user prompt (+ optional image reference):
  pptx      → PowerPoint presentations with styled slides
  xlsx      → Excel spreadsheets with formatting & charts
  docx      → Word documents with headings, tables, images
  pdf       → Styled PDF documents with cover page, sections, tables
  poster    → High-res poster images (PNG)
  infographic → Data-driven infographics (PNG)
  logo      → Branded logos (PNG)

Libraries: python-pptx, openpyxl, python-docx, fpdf2, Pillow, matplotlib
"""

import os
import re
import json
import math
import platform
import subprocess
from pathlib import Path
from datetime import datetime

# Valid image extensions for image_ref validation
_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".tif", ".webp", ".svg"}


def _extract_file_text(path: Path) -> str:
    """Extract readable text from non-text files (pdf, docx, pptx, xlsx)."""
    ext = path.suffix.lower()
    try:
        if ext == ".pdf":
            try:
                from fpdf import FPDF
                # fpdf2 is write-only, use PyPDF2 or pdfplumber if available
                try:
                    import pdfplumber
                    with pdfplumber.open(str(path)) as pdf:
                        return "\n".join(p.extract_text() or "" for p in pdf.pages)[:15000]
                except ImportError:
                    try:
                        from PyPDF2 import PdfReader
                        reader = PdfReader(str(path))
                        return "\n".join(p.extract_text() or "" for p in reader.pages)[:15000]
                    except ImportError:
                        pass
            except Exception:
                pass
            return ""
        elif ext == ".docx":
            try:
                from docx import Document
                doc = Document(str(path))
                return "\n".join(p.text for p in doc.paragraphs)[:15000]
            except Exception:
                return ""
        elif ext == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(str(path))
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            texts.append(shape.text_frame.text)
                return "\n".join(texts)[:15000]
            except Exception:
                return ""
        elif ext == ".xlsx":
            try:
                import openpyxl
                wb = openpyxl.load_workbook(str(path), data_only=True)
                texts = []
                for ws in wb.worksheets:
                    for row in ws.iter_rows(values_only=True):
                        texts.append(" | ".join(str(c) if c is not None else "" for c in row))
                return "\n".join(texts)[:15000]
            except Exception:
                return ""
    except Exception:
        pass
    return ""

# ─── Config & helpers ────────────────────────────────────────────────────────

# Fallback models tried in order if the primary model is unavailable
_FALLBACK_MODELS = ["gemini-2.5-flash", "gemini-2.5-flash-lite", "gemini-2.0-flash"]

def _gemini_generate(prompt: str, system: str = "", model: str = "gemini-2.5-flash") -> str:
    """Call Gemini and return plain text. Tries fallback models on transient errors."""
    from config.genai_client import generate_content, _is_transient_error

    cfg = {}
    if system:
        cfg["system_instruction"] = system

    # Build the model list: requested model first, then fallbacks
    models_to_try = [model] + [m for m in _FALLBACK_MODELS if m != model]

    last_exc = None
    for m in models_to_try:
        try:
            return generate_content(model=m, contents=prompt, config=cfg or None)
        except Exception as exc:
            last_exc = exc
            if _is_transient_error(exc) and m != models_to_try[-1]:
                print(f"[DocCreator] ⚠️  Model '{m}' unavailable ({exc}), trying fallback...")
                continue
            # Non-transient or last model — raise
            raise

def _save_dir() -> Path:
    d = Path.home() / "Desktop" / "JarvisMedia"
    d.mkdir(parents=True, exist_ok=True)
    return d

def _open_file(path: Path):
    """Auto-open a file on any OS."""
    try:
        if platform.system() == "Windows":
            os.startfile(str(path))
        elif platform.system() == "Darwin":
            subprocess.run(["open", str(path)])
        else:
            subprocess.run(["xdg-open", str(path)])
    except Exception:
        pass

def _hex_to_rgb(hex_color: str) -> tuple:
    """Convert #RRGGBB to (R, G, B)."""
    h = hex_color.lstrip("#")
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))


def _lighten(hex_color: str, amount: int = 40) -> tuple:
    """Lighten a hex color by `amount` per channel, capped at 255."""
    r, g, b = _hex_to_rgb(hex_color)
    return (min(255, r + amount), min(255, g + amount), min(255, b + amount))


def _darken(hex_color: str, amount: int = 40) -> tuple:
    """Darken a hex color by `amount` per channel, floored at 0."""
    r, g, b = _hex_to_rgb(hex_color)
    return (max(0, r - amount), max(0, g - amount), max(0, b - amount))


# ═══════════════════════════════════════════════════════════════════════════════
# 1. POWERPOINT (.pptx)
# ═══════════════════════════════════════════════════════════════════════════════

def _create_pptx(plan: dict, output_path: Path, image_ref: Path = None) -> str:
    """Build a PowerPoint from a Gemini-generated plan."""
    try:
        from pptx import Presentation
    except ImportError:
        return "python-pptx not installed. Run: py -3.11 -m pip install python-pptx"

    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides_data = plan.get("slides", [])
    color_scheme = plan.get("colors", {
        "bg": "#1a1a2e",
        "accent": "#e94560",
        "text": "#ffffff",
        "muted": "#aaaaaa"
    })

    bg_rgb   = RGBColor(*_hex_to_rgb(color_scheme.get("bg", "#1a1a2e")))
    acc_rgb  = RGBColor(*_hex_to_rgb(color_scheme.get("accent", "#e94560")))
    txt_rgb  = RGBColor(*_hex_to_rgb(color_scheme.get("text", "#ffffff")))
    muted_rgb = RGBColor(*_hex_to_rgb(color_scheme.get("muted", "#aaaaaa")))

    blank_layout = prs.slide_layouts[6]  # blank

    for idx, slide_d in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)

        # Background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = bg_rgb

        slide_type = slide_d.get("type", "content")

        if slide_type == "title":
            # Title slide — dynamic font size based on text length
            title_text = slide_d.get("title", "")
            sub_text = slide_d.get("subtitle", "")
            title_font_size = 54 if len(title_text) < 30 else 44 if len(title_text) < 60 else 36 if len(title_text) < 90 else 28
            sub_font_size = 24 if len(sub_text) < 50 else 20 if len(sub_text) < 80 else 16

            title_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(2.0), Inches(10), Inches(2.5)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title_text
            p.font.size = Pt(title_font_size)
            p.font.bold = True
            p.font.color.rgb = txt_rgb
            p.alignment = PP_ALIGN.CENTER

            sub_box = slide.shapes.add_textbox(
                Inches(2.5), Inches(4.5), Inches(8), Inches(1)
            )
            tf2 = sub_box.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = sub_text
            p2.font.size = Pt(sub_font_size)
            p2.font.color.rgb = muted_rgb
            p2.alignment = PP_ALIGN.CENTER

            # Accent bar — centered below title
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(5.5), Inches(4.3), Inches(2.3), Inches(0.08)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = acc_rgb
            bar.line.fill.background()

            # Decorative corner accent shapes
            corner = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0), Inches(0.4), Inches(7.5)
            )
            corner.fill.solid()
            corner.fill.fore_color.rgb = acc_rgb
            corner.line.fill.background()

        elif slide_type == "image":
            title_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(0.5), Inches(11), Inches(1)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = slide_d.get("title", "")
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = acc_rgb

            # Add image (from reference or a generated placeholder)
            img_source = image_ref if image_ref and image_ref.exists() else None
            if img_source and img_source.suffix.lower() in _IMAGE_EXTS:
                try:
                    slide.shapes.add_picture(
                        str(img_source), Inches(2.5), Inches(1.8), Inches(8), Inches(5)
                    )
                except Exception:
                    # Fallback: show a placeholder text instead of crashing
                    placeholder = slide.shapes.add_textbox(
                        Inches(2.5), Inches(3.5), Inches(8), Inches(2)
                    )
                    ptf = placeholder.text_frame
                    ptf.word_wrap = True
                    pp = ptf.paragraphs[0]
                    pp.text = f"[Image: {img_source.name}]"
                    pp.font.size = Pt(18)
                    pp.font.color.rgb = muted_rgb
                    pp.alignment = PP_ALIGN.CENTER

        elif slide_type == "bullets":
            title_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(0.5), Inches(11), Inches(1)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = slide_d.get("title", "")
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = acc_rgb

            bullets = slide_d.get("bullets", [])
            content_box = slide.shapes.add_textbox(
                Inches(1.2), Inches(1.8), Inches(10.5), Inches(5)
            )
            ctf = content_box.text_frame
            ctf.word_wrap = True

            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = ctf.paragraphs[0]
                else:
                    p = ctf.add_paragraph()
                p.text = f"▸  {bullet}"
                p.font.size = Pt(22)
                p.font.color.rgb = txt_rgb
                p.space_after = Pt(10)

        elif slide_type == "two_column":
            title_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(0.4), Inches(11), Inches(1)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = slide_d.get("title", "")
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = acc_rgb

            left_items  = slide_d.get("left", [])
            right_items = slide_d.get("right", [])

            for col_idx, (items, x_off) in enumerate([(left_items, 0.8), (right_items, 6.8)]):
                box = slide.shapes.add_textbox(
                    Inches(x_off), Inches(1.8), Inches(5.5), Inches(5)
                )
                btf = box.text_frame
                btf.word_wrap = True
                for i, item in enumerate(items):
                    if i == 0:
                        p = btf.paragraphs[0]
                    else:
                        p = btf.add_paragraph()
                    p.text = f"▸  {item}"
                    p.font.size = Pt(20)
                    p.font.color.rgb = txt_rgb
                    p.space_after = Pt(10)

        elif slide_type == "stats":
            title_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(0.4), Inches(11), Inches(1)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = slide_d.get("title", "")
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = acc_rgb

            stats = slide_d.get("stats", [])
            n = len(stats)
            if n > 0:
                card_w = min(3.5, 10.5 / n)
                gap = 0.4
                total_w = n * card_w + (n - 1) * gap
                start_x = (13.333 - total_w) / 2

                for i, stat in enumerate(stats):
                    x = start_x + i * (card_w + gap)
                    # Card background
                    card = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(x), Inches(2.2), Inches(card_w), Inches(3.8)
                    )
                    card.fill.solid()
                    card.fill.fore_color.rgb = RGBColor(
                        min(255, bg_rgb[0] + 25),
                        min(255, bg_rgb[1] + 25),
                        min(255, bg_rgb[2] + 25)
                    )
                    card.line.fill.background()

                    # Big number
                    num_box = slide.shapes.add_textbox(
                        Inches(x + 0.2), Inches(2.5), Inches(card_w - 0.4), Inches(1.5)
                    )
                    ntf = num_box.text_frame
                    np_ = ntf.paragraphs[0]
                    np_.text = stat.get("value", "0")
                    np_.font.size = Pt(44)
                    np_.font.bold = True
                    np_.font.color.rgb = acc_rgb
                    np_.alignment = PP_ALIGN.CENTER

                    # Label
                    lbl_box = slide.shapes.add_textbox(
                        Inches(x + 0.2), Inches(4.2), Inches(card_w - 0.4), Inches(1)
                    )
                    ltf = lbl_box.text_frame
                    ltf.word_wrap = True
                    lp = ltf.paragraphs[0]
                    lp.text = stat.get("label", "")
                    lp.font.size = Pt(16)
                    lp.font.color.rgb = muted_rgb
                    lp.alignment = PP_ALIGN.CENTER

        elif slide_type == "quote":
            # Quote slide — large italic quote with accent bar
            quote_text = slide_d.get("text", slide_d.get("title", ""))
            quote_author = slide_d.get("author", "")
            quote_font = 36 if len(quote_text) < 80 else 28 if len(quote_text) < 150 else 22

            # Large accent bar on the left
            accent_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(1.2), Inches(1.5), Inches(0.15), Inches(4.5)
            )
            accent_bar.fill.solid()
            accent_bar.fill.fore_color.rgb = acc_rgb
            accent_bar.line.fill.background()

            quote_box = slide.shapes.add_textbox(
                Inches(1.8), Inches(1.5), Inches(9.5), Inches(4)
            )
            qtf = quote_box.text_frame
            qtf.word_wrap = True
            qp = qtf.paragraphs[0]
            qp.text = f"\u201c{quote_text}\u201d"
            qp.font.size = Pt(quote_font)
            qp.font.italic = True
            qp.font.color.rgb = txt_rgb
            qp.alignment = PP_ALIGN.LEFT

            if quote_author:
                ap = qtf.add_paragraph()
                ap.text = f"\u2014 {quote_author}"
                ap.font.size = Pt(18)
                ap.font.color.rgb = muted_rgb
                ap.alignment = PP_ALIGN.LEFT
                ap.space_before = Pt(20)

        elif slide_type == "section_divider":
            # Section divider — centered title with accent stripe
            divider_title = slide_d.get("title", "")
            dt_font = 42 if len(divider_title) < 40 else 34
            section_num = slide_d.get("section_number", "")

            # Centered accent stripe
            stripe = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(5.0), Inches(3.2), Inches(3.3), Inches(0.1)
            )
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = acc_rgb
            stripe.line.fill.background()

            # Section number
            if section_num:
                num_box = slide.shapes.add_textbox(
                    Inches(2), Inches(2.2), Inches(9.3), Inches(0.8)
                )
                ntf = num_box.text_frame
                np = ntf.paragraphs[0]
                np.text = str(section_num)
                np.font.size = Pt(18)
                np.font.color.rgb = acc_rgb
                np.alignment = PP_ALIGN.CENTER

            title_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(3.5), Inches(10.3), Inches(1.5)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = divider_title
            p.font.size = Pt(dt_font)
            p.font.bold = True
            p.font.color.rgb = txt_rgb
            p.alignment = PP_ALIGN.CENTER

            subtitle = slide_d.get("subtitle", "")
            if subtitle:
                sub_box = slide.shapes.add_textbox(
                    Inches(2.5), Inches(5.0), Inches(8.3), Inches(0.8)
                )
                stf = sub_box.text_frame
                stf.word_wrap = True
                sp = stf.paragraphs[0]
                sp.text = subtitle
                sp.font.size = Pt(18)
                sp.font.color.rgb = muted_rgb
                sp.alignment = PP_ALIGN.CENTER

        else:
            # Generic content slide — with accent bar decoration
            content_title = slide_d.get("title", "Slide")
            title_font_size = 36 if len(content_title) < 40 else 30

            title_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(0.5), Inches(11), Inches(1)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = content_title
            p.font.size = Pt(title_font_size)
            p.font.bold = True
            p.font.color.rgb = acc_rgb

            # Accent underline
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.8), Inches(1.3), Inches(1.8), Inches(0.06)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = acc_rgb
            bar.line.fill.background()

            body_text = slide_d.get("body", "")
            if body_text:
                body_box = slide.shapes.add_textbox(
                    Inches(1.2), Inches(1.6), Inches(10.5), Inches(5.4)
                )
                btf = body_box.text_frame
                btf.word_wrap = True
                bp = btf.paragraphs[0]
                bp.text = body_text
                bp.font.size = Pt(20)
                bp.font.color.rgb = txt_rgb

    prs.save(str(output_path))
    return f"Presentation saved: {output_path.name} ({len(slides_data)} slides)"


# ═══════════════════════════════════════════════════════════════════════════════
# 2. EXCEL SPREADSHEET (.xlsx)
# ═══════════════════════════════════════════════════════════════════════════════

def _create_xlsx(plan: dict, output_path: Path) -> str:
    """Build an Excel spreadsheet from a Gemini-generated plan."""
    try:
        import openpyxl
    except ImportError:
        return "openpyxl not installed. Run: py -3.11 -m pip install openpyxl"

    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.chart import BarChart, PieChart, LineChart, Reference
    from openpyxl.utils import get_column_letter

    wb = openpyxl.Workbook()

    color_scheme = plan.get("colors", {
        "header_bg": "#1a1a2e",
        "header_text": "#ffffff",
        "accent": "#e94560",
        "alt_row": "#f0f0f5"
    })

    header_fill  = PatternFill("solid", fgColor=color_scheme.get("header_bg", "#1a1a2e").lstrip("#"))
    header_font  = Font(bold=True, color=color_scheme.get("header_text", "#ffffff").lstrip("#"), size=12)
    accent_fill  = PatternFill("solid", fgColor=color_scheme.get("accent", "#e94560").lstrip("#"))
    alt_fill     = PatternFill("solid", fgColor=color_scheme.get("alt_row", "#f0f0f5").lstrip("#"))
    thin_border  = Border(
        left=Side(style="thin"), right=Side(style="thin"),
        top=Side(style="thin"), bottom=Side(style="thin")
    )
    center_align = Alignment(horizontal="center", vertical="center")

    sheets = plan.get("sheets", [{"name": "Data", "headers": [], "rows": []}])

    for sheet_d in sheets:
        name = sheet_d.get("name", "Sheet1")[:31]  # Excel limit
        wb.create_sheet(title=name)
        ws = wb[name]

        headers = sheet_d.get("headers", [])
        rows    = sheet_d.get("rows", [])

        # Write headers
        for col_idx, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col_idx, value=header)
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = center_align
            cell.border = thin_border

        # Write data rows
        for row_idx, row_data in enumerate(rows, 2):
            for col_idx, value in enumerate(row_data, 1):
                cell = ws.cell(row=row_idx, column=col_idx, value=value)
                cell.border = thin_border
                cell.alignment = Alignment(horizontal="center", vertical="center")
                if row_idx % 2 == 0:
                    cell.fill = alt_fill

        # Auto-fit column widths
        for col_idx in range(1, len(headers) + 1):
            max_len = max(
                len(str(ws.cell(row=r, column=col_idx).value or ""))
                for r in range(1, len(rows) + 2)
            )
            ws.column_dimensions[get_column_letter(col_idx)].width = max(max_len + 4, 12)

        # Add chart if plan specifies one
        chart_spec = sheet_d.get("chart")
        if chart_spec and rows:
                chart_type = chart_spec.get("type", "bar")
                data_rows = len(rows) + 1
                data_cols = len(headers)

                if chart_type == "bar":
                    chart = BarChart()
                elif chart_type == "pie":
                    chart = PieChart()
                elif chart_type == "line":
                    chart = LineChart()
                else:
                    chart = BarChart()

                chart.title = chart_spec.get("title", "Chart")
                chart.style = 10
                chart.width  = 20
                chart.height = 12

                # Validate column indices to prevent crashes
                cat_col = max(1, min(chart_spec.get("category_column", 1), data_cols))
                val_col = max(1, min(chart_spec.get("value_column", 2), data_cols))
                if cat_col == val_col:
                    val_col = min(cat_col + 1, data_cols) if data_cols > 1 else 1

                data = Reference(ws, min_col=val_col, min_row=1, max_row=data_rows, max_col=val_col)
                cats = Reference(ws, min_col=cat_col, min_row=2, max_row=data_rows)

                chart.add_data(data, titles_from_data=True)
                chart.set_categories(cats)
                chart.shape = 4

                chart_row = data_rows + 3
                ws.add_chart(chart, f"A{chart_row}")

    # Remove default sheet
    if "Sheet" in wb.sheetnames:
        del wb["Sheet"]

    wb.save(str(output_path))
    return f"Spreadsheet saved: {output_path.name} ({len(sheets)} sheets)"


# ═══════════════════════════════════════════════════════════════════════════════
# 3. WORD DOCUMENT (.docx)
# ═══════════════════════════════════════════════════════════════════════════════

def _create_docx(plan: dict, output_path: Path, image_ref: Path = None) -> str:
    """Build a Word document from a Gemini-generated plan."""
    try:
        from docx import Document
    except ImportError:
        return "python-docx not installed. Run: py -3.11 -m pip install python-docx"

    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT

    doc = Document()

    # Page margins
    for section in doc.sections:
        section.top_margin    = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin   = Inches(1.2)
        section.right_margin  = Inches(1.2)

    color_scheme = plan.get("colors", {
        "heading": "#1a1a2e",
        "accent": "#e94560",
        "body": "#333333"
    })

    heading_rgb = RGBColor(*_hex_to_rgb(color_scheme.get("heading", "#1a1a2e")))
    accent_rgb  = RGBColor(*_hex_to_rgb(color_scheme.get("accent", "#e94560")))

    elements = plan.get("elements", [])

    for el in elements:
        el_type = el.get("type", "paragraph")

        if el_type == "heading":
            level = min(el.get("level", 1), 4)
            heading = doc.add_heading(el.get("text", ""), level=level)
            for run in heading.runs:
                run.font.color.rgb = heading_rgb

        elif el_type == "paragraph":
            p = doc.add_paragraph(el.get("text", ""))
            p.style = doc.styles["Normal"]
            for run in p.runs:
                run.font.size = Pt(11)
                run.font.color.rgb = RGBColor(*_hex_to_rgb(color_scheme.get("body", "#333333")))

        elif el_type == "bullet":
            p = doc.add_paragraph(el.get("text", ""), style="List Bullet")
            for run in p.runs:
                run.font.size = Pt(11)

        elif el_type == "numbered":
            p = doc.add_paragraph(el.get("text", ""), style="List Number")
            for run in p.runs:
                run.font.size = Pt(11)

        elif el_type == "table":
            headers = el.get("headers", [])
            rows    = el.get("rows", [])
            if headers:
                table = doc.add_table(rows=1 + len(rows), cols=len(headers))
                table.style = "Light Grid Accent 1"
                table.alignment = WD_TABLE_ALIGNMENT.CENTER

                # Header row
                for i, h in enumerate(headers):
                    cell = table.rows[0].cells[i]
                    cell.text = h
                    for p in cell.paragraphs:
                        for run in p.runs:
                            run.font.bold = True
                            run.font.size = Pt(10)

                # Data rows
                for r_idx, row_data in enumerate(rows):
                    for c_idx, val in enumerate(row_data):
                        cell = table.rows[r_idx + 1].cells[c_idx]
                        cell.text = str(val)
                        for p in cell.paragraphs:
                            for run in p.runs:
                                run.font.size = Pt(10)

        elif el_type == "image":
            caption = el.get("caption", "")
            img_source = image_ref if image_ref and image_ref.exists() else None
            if img_source and img_source.suffix.lower() in _IMAGE_EXTS:
                try:
                    doc.add_picture(str(img_source), width=Inches(5))
                    last_paragraph = doc.paragraphs[-1]
                    last_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                except Exception:
                    p = doc.add_paragraph(f"[Image: {img_source.name}]")
                    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    for run in p.runs:
                        run.font.italic = True
                if caption:
                    cap_p = doc.add_paragraph(caption)
                    cap_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    for run in cap_p.runs:
                        run.font.size = Pt(9)
                        run.font.italic = True

        elif el_type == "divider":
            doc.add_paragraph("─" * 60)

        elif el_type == "quote":
            # Styled quote with accent left border + shading
            p = doc.add_paragraph(el.get("text", ""))
            p.style = doc.styles["Normal"]
            pf = p.paragraph_format
            pf.left_indent = Inches(0.8)
            pf.right_indent = Inches(0.4)
            pf.space_before = Pt(12)
            pf.space_after = Pt(12)
            # Add a left border for visual accent
            from docx.oxml.ns import qn
            from docx.oxml import OxmlElement
            pBdr = OxmlElement('w:pBdr')
            left = OxmlElement('w:left')
            left.set(qn('w:val'), 'single')
            left.set(qn('w:sz'), '24')
            left.set(qn('w:space'), '8')
            left.set(qn('w:color'), color_scheme.get("accent", "#e94560").lstrip('#'))
            pBdr.append(left)
            p._p.get_or_add_pPr().append(pBdr)
            for run in p.runs:
                run.font.italic = True
                run.font.size = Pt(12)
                run.font.color.rgb = accent_rgb

    doc.save(str(output_path))
    return f"Document saved: {output_path.name}"


# ═══════════════════════════════════════════════════════════════════════════════
# 4. POSTER (PNG via Pillow)
# ═══════════════════════════════════════════════════════════════════════════════

def _create_poster(plan: dict, output_path: Path, image_ref: Path = None) -> str:
    """Build a high-res poster image from a plan."""
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        return "Pillow not installed. Run: py -3.11 -m pip install Pillow"

    from PIL import Image, ImageDraw, ImageFont

    width  = plan.get("width", 2480)
    height = plan.get("height", 3508)  # A4 at 300 DPI
    color_scheme = plan.get("colors", {
        "bg": "#1a1a2e",
        "accent": "#e94560",
        "text": "#ffffff",
        "muted": "#aaaaaa"
    })

    img = Image.new("RGB", (width, height), color_scheme.get("bg", "#1a1a2e"))
    draw = ImageDraw.Draw(img)

    bg_rgb   = _hex_to_rgb(color_scheme.get("bg", "#1a1a2e"))
    acc_rgb  = _hex_to_rgb(color_scheme.get("accent", "#e94560"))
    txt_rgb  = _hex_to_rgb(color_scheme.get("text", "#ffffff"))
    muted_rgb = _hex_to_rgb(color_scheme.get("muted", "#aaaaaa"))

    # Try to load a nice font, fallback to default
    def _get_font(size, bold=False):
        font_paths = [
            "C:/Windows/Fonts/arialbd.ttf" if bold else "C:/Windows/Fonts/arial.ttf",
            "C:/Windows/Fonts/segoeui.ttf",
            "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        ]
        for fp in font_paths:
            if os.path.exists(fp):
                try:
                    return ImageFont.truetype(fp, size)
                except Exception:
                    continue
        return ImageFont.load_default()

    elements = plan.get("elements", [])

    for el in elements:
        el_type = el.get("type", "text")

        if el_type == "title":
            size = el.get("font_size", 120)
            font = _get_font(size, bold=True)
            text = el.get("text", "")
            x = el.get("x", width // 2)
            y = el.get("y", 200)
            color = _hex_to_rgb(el.get("color", color_scheme.get("text", "#ffffff")))
            # Center text
            bbox = draw.textbbox((0, 0), text, font=font)
            tw = bbox[2] - bbox[0]
            draw.text((x - tw // 2, y), text, fill=color, font=font)

        elif el_type == "subtitle":
            size = el.get("font_size", 48)
            font = _get_font(size, bold=False)
            text = el.get("text", "")
            x = el.get("x", width // 2)
            y = el.get("y", 350)
            color = _hex_to_rgb(el.get("color", color_scheme.get("muted", "#aaaaaa")))
            bbox = draw.textbbox((0, 0), text, font=font)
            tw = bbox[2] - bbox[0]
            draw.text((x - tw // 2, y), text, fill=color, font=font)

        elif el_type == "text":
            size = el.get("font_size", 36)
            font = _get_font(size)
            text = el.get("text", "")
            x = el.get("x", 200)
            y = el.get("y", 500)
            color = _hex_to_rgb(el.get("color", color_scheme.get("text", "#ffffff")))
            max_width = el.get("max_width", width - 400)

            # Word wrap
            words = text.split()
            lines = []
            current_line = ""
            for word in words:
                test_line = f"{current_line} {word}".strip()
                bbox = draw.textbbox((0, 0), test_line, font=font)
                if bbox[2] - bbox[0] <= max_width:
                    current_line = test_line
                else:
                    if current_line:
                        lines.append(current_line)
                    current_line = word
            if current_line:
                lines.append(current_line)

            for i, line in enumerate(lines):
                draw.text((x, y + i * (size + 10)), line, fill=color, font=font)

        elif el_type == "rectangle":
            x1 = el.get("x", 100)
            y1 = el.get("y", 100)
            w  = el.get("width", 400)
            h  = el.get("height", 300)
            color = _hex_to_rgb(el.get("color", color_scheme.get("accent", "#e94560")))
            radius = el.get("radius", 0)
            if radius > 0:
                draw.rounded_rectangle([x1, y1, x1 + w, y1 + h], radius=radius, fill=color)
            else:
                draw.rectangle([x1, y1, x1 + w, y1 + h], fill=color)

        elif el_type == "circle":
            cx = el.get("x", 500)
            cy = el.get("y", 500)
            r  = el.get("radius", 100)
            color = _hex_to_rgb(el.get("color", color_scheme.get("accent", "#e94560")))
            draw.ellipse([cx - r, cy - r, cx + r, cy + r], fill=color)

        elif el_type == "line":
            x1, y1 = el.get("x1", 0), el.get("y1", 0)
            x2, y2 = el.get("x2", 100), el.get("y2", 100)
            color = _hex_to_rgb(el.get("color", color_scheme.get("accent", "#e94560")))
            thickness = el.get("thickness", 4)
            draw.line([x1, y1, x2, y2], fill=color, width=thickness)

        elif el_type == "image":
            img_x = el.get("x", 200)
            img_y = el.get("y", 800)
            img_w = el.get("width", 800)
            img_h = el.get("height", 600)
            source = image_ref if image_ref and image_ref.exists() else None
            if source and source.suffix.lower() in _IMAGE_EXTS:
                try:
                    ref_img = Image.open(source).resize((img_w, img_h), Image.LANCZOS)
                    # RGBA-safe paste: use mask for transparency
                    if ref_img.mode == "RGBA":
                        img.paste(ref_img, (img_x, img_y), ref_img)
                    else:
                        img.paste(ref_img, (img_x, img_y))
                except Exception:
                    pass

    img.save(str(output_path), "PNG", quality=95)
    return f"Poster saved: {output_path.name} ({width}x{height}px)"


# ═══════════════════════════════════════════════════════════════════════════════
# 5. INFOGRAPHIC (PNG via Pillow + matplotlib)
# ═══════════════════════════════════════════════════════════════════════════════

def _create_infographic(plan: dict, output_path: Path) -> str:
    """Build a data-driven infographic combining Pillow canvas with matplotlib charts."""
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        return "Pillow not installed. Run: py -3.11 -m pip install Pillow"
    try:
        import matplotlib
    except ImportError:
        return "matplotlib not installed. Run: py -3.11 -m pip install matplotlib"

    from PIL import Image, ImageDraw, ImageFont
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    from io import BytesIO

    width  = plan.get("width", 1200)
    height = plan.get("height", 2400)
    color_scheme = plan.get("colors", {
        "bg": "#0f0f23",
        "accent": "#e94560",
        "secondary": "#00e5ff",
        "text": "#ffffff",
        "muted": "#888899"
    })

    img = Image.new("RGB", (width, height), color_scheme.get("bg", "#0f0f23"))
    draw = ImageDraw.Draw(img)

    acc_rgb  = _hex_to_rgb(color_scheme.get("accent", "#e94560"))
    sec_rgb  = _hex_to_rgb(color_scheme.get("secondary", "#00e5ff"))
    txt_rgb  = _hex_to_rgb(color_scheme.get("text", "#ffffff"))
    muted_rgb = _hex_to_rgb(color_scheme.get("muted", "#888899"))

    def _get_font(size, bold=False):
        font_paths = [
            "C:/Windows/Fonts/arialbd.ttf" if bold else "C:/Windows/Fonts/arial.ttf",
            "C:/Windows/Fonts/segoeui.ttf",
        ]
        for fp in font_paths:
            if os.path.exists(fp):
                try:
                    return ImageFont.truetype(fp, size)
                except Exception:
                    continue
        return ImageFont.load_default()

    # Render matplotlib charts as images and paste onto canvas
    charts = plan.get("charts", [])
    for chart_d in charts:
        chart_type = chart_d.get("type", "bar")
        chart_data = chart_d.get("data", {})
        chart_x    = chart_d.get("x", 100)
        chart_y    = chart_d.get("y", 400)
        chart_w    = chart_d.get("width", 500)
        chart_h    = chart_d.get("height", 350)

        fig, ax = plt.subplots(figsize=(chart_w / 100, chart_h / 100), dpi=100)
        fig.patch.set_facecolor(color_scheme.get("bg", "#0f0f23"))
        ax.set_facecolor(color_scheme.get("bg", "#0f0f23"))

        labels = chart_data.get("labels", [])
        values = chart_data.get("values", [])
        chart_title = chart_d.get("title", "")

        if chart_type == "bar":
            colors = [f"#{color_scheme.get('accent', '#e94560').lstrip('#')}"] * len(values)
            ax.bar(labels, values, color=colors, width=0.6)
        elif chart_type == "pie":
            colors_pie = [
                color_scheme.get("accent", "#e94560"),
                color_scheme.get("secondary", "#00e5ff"),
                "#ffd700", "#ff6b35", "#7b2ff7", "#00cc88"
            ]
            ax.pie(values, labels=labels, colors=[c.lstrip("#") for c in colors_pie[:len(values)]],
                   autopct="%1.0f%%", textprops={"color": "white", "fontsize": 10})
        elif chart_type == "line":
            ax.plot(labels, values, color=f"#{color_scheme.get('secondary', '#00e5ff').lstrip('#')}",
                    linewidth=3, marker="o", markersize=8)
            ax.fill_between(range(len(labels)), values, alpha=0.15,
                           color=f"#{color_scheme.get('secondary', '#00e5ff').lstrip('#')}")

        if chart_title:
            ax.set_title(chart_title, color="white", fontsize=14, fontweight="bold", pad=10)
        ax.tick_params(colors="white", labelsize=9)
        for spine in ax.spines.values():
            spine.set_color("#333355")

        buf = BytesIO()
        fig.savefig(buf, format="png", bbox_inches="tight", facecolor=fig.get_facecolor(), dpi=100)
        plt.close(fig)
        buf.seek(0)

        chart_img = Image.open(buf).resize((chart_w, chart_h), Image.LANCZOS)
        if chart_img.mode == "RGBA":
            img.paste(chart_img, (chart_x, chart_y), chart_img)
        else:
            img.paste(chart_img, (chart_x, chart_y))

    # Render text elements
    elements = plan.get("elements", [])
    for el in elements:
        el_type = el.get("type", "text")

        if el_type == "title":
            size = el.get("font_size", 64)
            font = _get_font(size, bold=True)
            text = el.get("text", "")
            x = el.get("x", width // 2)
            y = el.get("y", 80)
            color = _hex_to_rgb(el.get("color", color_scheme.get("text", "#ffffff")))
            bbox = draw.textbbox((0, 0), text, font=font)
            tw = bbox[2] - bbox[0]
            draw.text((x - tw // 2, y), text, fill=color, font=font)

        elif el_type == "stat_card":
            card_x = el.get("x", 100)
            card_y = el.get("y", 300)
            card_w = el.get("width", 250)
            card_h = el.get("height", 180)

            # Card background — light tint of accent (RGB-safe)
            card_bg = _lighten(color_scheme.get("accent", "#e94560"), 60)
            draw.rounded_rectangle(
                [card_x, card_y, card_x + card_w, card_y + card_h],
                radius=16, fill=card_bg
            )

            # Value
            val_font = _get_font(42, bold=True)
            val_text = el.get("value", "0")
            bbox = draw.textbbox((0, 0), val_text, font=val_font)
            tw = bbox[2] - bbox[0]
            draw.text(
                (card_x + card_w // 2 - tw // 2, card_y + 30),
                val_text, fill=acc_rgb, font=val_font
            )

            # Label
            lbl_font = _get_font(18)
            lbl_text = el.get("label", "")
            bbox = draw.textbbox((0, 0), lbl_text, font=lbl_font)
            tw = bbox[2] - bbox[0]
            draw.text(
                (card_x + card_w // 2 - tw // 2, card_y + 90),
                lbl_text, fill=muted_rgb, font=lbl_font
            )

        elif el_type == "text":
            size = el.get("font_size", 28)
            font = _get_font(size)
            text = el.get("text", "")
            x = el.get("x", 100)
            y = el.get("y", 600)
            color = _hex_to_rgb(el.get("color", color_scheme.get("text", "#ffffff")))
            max_width = el.get("max_width", width - 200)

            # Word wrap for infographic text
            words = text.split()
            lines = []
            current_line = ""
            for word in words:
                test_line = f"{current_line} {word}".strip()
                bbox = draw.textbbox((0, 0), test_line, font=font)
                if bbox[2] - bbox[0] <= max_width:
                    current_line = test_line
                else:
                    if current_line:
                        lines.append(current_line)
                    current_line = word
            if current_line:
                lines.append(current_line)

            for i, line in enumerate(lines):
                draw.text((x, y + i * (size + 8)), line, fill=color, font=font)

        elif el_type == "divider":
            y = el.get("y", 380)
            draw.line([60, y, width - 60, y], fill=acc_rgb, width=2)

    img.save(str(output_path), "PNG", quality=95)
    return f"Infographic saved: {output_path.name} ({width}x{height}px)"


# ═══════════════════════════════════════════════════════════════════════════════
# 6. LOGO (PNG via Pillow)
# ═══════════════════════════════════════════════════════════════════════════════

def _create_logo(plan: dict, output_path: Path) -> str:
    """Build a branded logo image."""
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        return "Pillow not installed. Run: py -3.11 -m pip install Pillow"

    from PIL import Image, ImageDraw, ImageFont

    size = plan.get("size", 1024)
    color_scheme = plan.get("colors", {
        "bg": "#1a1a2e",
        "accent": "#e94560",
        "text": "#ffffff"
    })

    bg_mode = "RGBA" if plan.get("transparent_bg", False) else "RGB"
    bg_color = (0, 0, 0, 0) if bg_mode == "RGBA" else color_scheme.get("bg", "#1a1a2e")

    img = Image.new(bg_mode, (size, size), bg_color)
    draw = ImageDraw.Draw(img)

    bg_rgb   = _hex_to_rgb(color_scheme.get("bg", "#1a1a2e"))
    acc_rgb  = _hex_to_rgb(color_scheme.get("accent", "#e94560"))
    txt_rgb  = _hex_to_rgb(color_scheme.get("text", "#ffffff"))

    def _get_font(fsize, bold=True):
        font_paths = [
            "C:/Windows/Fonts/arialbd.ttf" if bold else "C:/Windows/Fonts/arial.ttf",
            "C:/Windows/Fonts/segoeui.ttf",
        ]
        for fp in font_paths:
            if os.path.exists(fp):
                try:
                    return ImageFont.truetype(fp, fsize)
                except Exception:
                    continue
        return ImageFont.load_default()

    elements = plan.get("elements", [])

    for el in elements:
        el_type = el.get("type", "text")

        if el_type == "shape":
            shape = el.get("shape", "circle")
            cx = el.get("x", size // 2)
            cy = el.get("y", size // 2)
            r  = el.get("radius", size // 3)
            color = _hex_to_rgb(el.get("color", color_scheme.get("accent", "#e94560")))

            if shape == "circle":
                draw.ellipse([cx - r, cy - r, cx + r, cy + r], fill=color)
            elif shape == "rounded_rect":
                w = el.get("width", size // 2)
                h = el.get("height", size // 2)
                rad = el.get("corner_radius", 30)
                draw.rounded_rectangle(
                    [cx - w // 2, cy - h // 2, cx + w // 2, cy + h // 2],
                    radius=rad, fill=color
                )
            elif shape == "hexagon":
                points = []
                for i in range(6):
                    angle = math.radians(60 * i - 30)
                    px = cx + r * math.cos(angle)
                    py = cy + r * math.sin(angle)
                    points.append((px, py))
                draw.polygon(points, fill=color)
            elif shape == "diamond":
                points = [(cx, cy - r), (cx + r, cy), (cx, cy + r), (cx - r, cy)]
                draw.polygon(points, fill=color)
            elif shape == "star":
                points = []
                for i in range(10):
                    angle = math.radians(36 * i - 90)
                    rad = r if i % 2 == 0 else r * 0.4
                    px = cx + rad * math.cos(angle)
                    py = cy + rad * math.sin(angle)
                    points.append((px, py))
                draw.polygon(points, fill=color)

        elif el_type == "initials":
            text = el.get("text", "J")
            fsize = el.get("font_size", size // 3)
            font = _get_font(fsize, bold=True)
            color = _hex_to_rgb(el.get("color", color_scheme.get("text", "#ffffff")))
            x = el.get("x", size // 2)
            y = el.get("y", size // 2)
            bbox = draw.textbbox((0, 0), text, font=font)
            tw = bbox[2] - bbox[0]
            th = bbox[3] - bbox[1]
            draw.text((x - tw // 2, y - th // 2 - 10), text, fill=color, font=font)

        elif el_type == "text":
            text = el.get("text", "")
            fsize = el.get("font_size", size // 10)
            font = _get_font(fsize, bold=True)
            color = _hex_to_rgb(el.get("color", color_scheme.get("text", "#ffffff")))
            x = el.get("x", size // 2)
            y = el.get("y", size * 3 // 4)
            bbox = draw.textbbox((0, 0), text, font=font)
            tw = bbox[2] - bbox[0]
            # Scale down font if text is too wide for the logo
            if tw > size * 0.85:
                new_size = int(fsize * size * 0.85 / tw)
                font = _get_font(new_size, bold=True)
                bbox = draw.textbbox((0, 0), text, font=font)
                tw = bbox[2] - bbox[0]
            draw.text((x - tw // 2, y), text, fill=color, font=font)

        elif el_type == "ring":
            cx = el.get("x", size // 2)
            cy = el.get("y", size // 2)
            r  = el.get("radius", size // 3)
            width_line = el.get("thickness", 8)
            color = _hex_to_rgb(el.get("color", color_scheme.get("accent", "#e94560")))
            draw.ellipse(
                [cx - r, cy - r, cx + r, cy + r],
                outline=color, width=width_line
            )

    img.save(str(output_path), "PNG")
    return f"Logo saved: {output_path.name} ({size}x{size}px)"


# ═══════════════════════════════════════════════════════════════════════════════
# 7. PDF DOCUMENT (.pdf)
# ═══════════════════════════════════════════════════════════════════════════════

def _create_pdf(plan: dict, output_path: Path, image_ref: Path = None) -> str:
    """Build a styled PDF document from a Gemini-generated plan."""
    try:
        from fpdf import FPDF
    except ImportError:
        return "fpdf2 not installed. Run: py -3.11 -m pip install fpdf2"

    from fpdf import FPDF

    color_scheme = plan.get("colors", {
        "bg": "#ffffff",
        "accent": "#1a1a2e",
        "text": "#333333",
        "muted": "#888888"
    })

    bg_rgb   = _hex_to_rgb(color_scheme.get("bg", "#ffffff"))
    acc_rgb  = _hex_to_rgb(color_scheme.get("accent", "#1a1a2e"))
    txt_rgb  = _hex_to_rgb(color_scheme.get("text", "#333333"))
    muted_rgb = _hex_to_rgb(color_scheme.get("muted", "#888888"))

    # Try to find a TTF font on the system for better rendering
    def _find_ttf(bold=False):
        candidates = [
            "C:/Windows/Fonts/arialbd.ttf" if bold else "C:/Windows/Fonts/arial.ttf",
            "C:/Windows/Fonts/segoeuib.ttf" if bold else "C:/Windows/Fonts/segoeui.ttf",
            "C:/Windows/Fonts/calibrib.ttf" if bold else "C:/Windows/Fonts/calibri.ttf",
        ]
        for fp in candidates:
            if os.path.exists(fp):
                return fp
        return None

    font_regular = _find_ttf(bold=False)
    font_bold    = _find_ttf(bold=True)

    class StyledPDF(FPDF):
        def __init__(self):
            super().__init__()
            self._custom_header = plan.get("header", "")
            self._custom_footer = plan.get("footer", "")

        def header(self):
            if self._custom_header:
                self.set_font("Helvetica", "B", 9)
                self.set_text_color(*muted_rgb)
                self.cell(0, 8, self._custom_header, align="R")
                self.ln(12)

        def footer(self):
            self.set_y(-15)
            self.set_font("Helvetica", "I", 8)
            self.set_text_color(*muted_rgb)
            label = self._custom_footer or f"Page {self.page_no()}/{{nb}}"
            self.cell(0, 10, label, align="C")

    pdf = StyledPDF()
    pdf.alias_nb_pages()
    pdf.set_auto_page_break(auto=True, margin=20)

    # Register custom fonts if available (need BOTH regular + bold for bold styling)
    font_family = "Helvetica"
    if font_regular:
        try:
            pdf.add_font("Custom", "", font_regular)
            if font_bold:
                pdf.add_font("Custom", "B", font_bold)
                font_family = "Custom"
            # If only regular found, stay with Helvetica (bold calls would crash)
        except Exception:
            font_family = "Helvetica"

    # ── Cover page ────────────────────────────────────────────────────────
    cover = plan.get("cover", {})
    if cover:
        pdf.add_page()
        # Background
        pdf.set_fill_color(*bg_rgb)
        pdf.rect(0, 0, 210, 297, "F")

        # Accent stripe
        stripe_color = _hex_to_rgb(cover.get("stripe_color", color_scheme.get("accent", "#1a1a2e")))
        pdf.set_fill_color(*stripe_color)
        pdf.rect(0, 100, 210, 6, "F")

        # Title
        pdf.set_y(120)
        pdf.set_font(font_family, "B", 32)
        pdf.set_text_color(*acc_rgb)
        title = cover.get("title", "")
        pdf.multi_cell(0, 14, title, align="C")

        # Subtitle
        pdf.ln(6)
        pdf.set_font(font_family, "", 16)
        pdf.set_text_color(*muted_rgb)
        subtitle = cover.get("subtitle", "")
        if subtitle:
            pdf.multi_cell(0, 10, subtitle, align="C")

        # Date / Author
        pdf.ln(20)
        pdf.set_font(font_family, "", 11)
        pdf.set_text_color(*muted_rgb)
        meta_parts = []
        if cover.get("author"):
            meta_parts.append(f"By {cover['author']}")
        if cover.get("date"):
            meta_parts.append(cover["date"])
        elif plan.get("date"):
            meta_parts.append(plan["date"])
        if meta_parts:
            pdf.cell(0, 8, " | ".join(meta_parts), align="C")

    # ── Body pages ────────────────────────────────────────────────────────
    sections = plan.get("sections", [])

    for section in sections:
        # Check if we need a new page
        if pdf.get_y() > 240:
            pdf.add_page()

        sec_type = section.get("type", "text")

        if sec_type == "heading":
            level = section.get("level", 1)
            pdf.ln(6)
            if level == 1:
                pdf.set_font(font_family, "B", 22)
                pdf.set_text_color(*acc_rgb)
                pdf.multi_cell(0, 12, section.get("text", ""))
                # Accent underline
                y = pdf.get_y()
                pdf.set_fill_color(*acc_rgb)
                pdf.rect(10, y + 1, 50, 1.5, "F")
                pdf.ln(6)
            elif level == 2:
                pdf.set_font(font_family, "B", 17)
                pdf.set_text_color(*acc_rgb)
                pdf.multi_cell(0, 10, section.get("text", ""))
                pdf.ln(3)
            else:
                pdf.set_font(font_family, "B", 13)
                pdf.set_text_color(*txt_rgb)
                pdf.multi_cell(0, 8, section.get("text", ""))
                pdf.ln(2)

        elif sec_type == "text":
            pdf.set_font(font_family, "", 11)
            pdf.set_text_color(*txt_rgb)
            pdf.multi_cell(0, 7, section.get("text", ""))
            pdf.ln(4)

        elif sec_type == "bullet":
            pdf.set_font(font_family, "", 11)
            pdf.set_text_color(*txt_rgb)
            items = section.get("items", [])
            if not items and section.get("text"):
                items = [section["text"]]
            for item in items:
                bullet = chr(8226) if font_family != "Helvetica" else "-"
                pdf.cell(8, 7, bullet)
                pdf.multi_cell(0, 7, item)
                pdf.ln(1)
            pdf.ln(3)

        elif sec_type == "numbered":
            pdf.set_font(font_family, "", 11)
            pdf.set_text_color(*txt_rgb)
            items = section.get("items", [])
            if not items and section.get("text"):
                items = [section["text"]]
            for i, item in enumerate(items, 1):
                pdf.cell(10, 7, f"{i}.")
                pdf.multi_cell(0, 7, item)
                pdf.ln(1)
            pdf.ln(3)

        elif sec_type == "quote":
            pdf.set_font(font_family, "I", 12)
            pdf.set_text_color(*acc_rgb)
            pdf.set_x(20)
            pdf.multi_cell(170, 7, section.get("text", ""))
            pdf.set_text_color(*txt_rgb)
            pdf.ln(4)

        elif sec_type == "table":
            headers = section.get("headers", [])
            rows    = section.get("rows", [])
            if headers:
                col_w = 190 / len(headers)
                # Header row
                pdf.set_font(font_family, "B", 10)
                pdf.set_fill_color(*acc_rgb)
                pdf.set_text_color(255, 255, 255)
                for h in headers:
                    pdf.cell(col_w, 9, str(h), border=1, align="C", fill=True)
                pdf.ln()
                # Data rows
                pdf.set_font(font_family, "", 9)
                pdf.set_text_color(*txt_rgb)
                for r_idx, row in enumerate(rows):
                    if r_idx % 2 == 0:
                        pdf.set_fill_color(240, 240, 245)
                        fill = True
                    else:
                        fill = False
                    for val in row:
                        pdf.cell(col_w, 8, str(val), border=1, align="C", fill=fill)
                    pdf.ln()
                pdf.ln(5)

        elif sec_type == "divider":
            y = pdf.get_y()
            pdf.set_draw_color(*muted_rgb)
            pdf.set_line_width(0.3)
            pdf.line(20, y, 190, y)
            pdf.ln(6)

        elif sec_type == "image":
            img_source = image_ref if image_ref and image_ref.exists() else None
            img_path = section.get("path")
            if img_path and os.path.exists(img_path):
                img_source = Path(img_path)
            if img_source:
                try:
                    pdf.image(str(img_source), x=20, w=170)
                    pdf.ln(5)
                except Exception:
                    pdf.set_font(font_family, "I", 10)
                    pdf.set_text_color(*muted_rgb)
                    pdf.cell(0, 8, "[Image could not be embedded]", align="C")
                    pdf.ln(10)

        elif sec_type == "stat_card":
            # Render stat cards as a horizontal row of boxes
            cards = section.get("cards", [])
            if not cards and section.get("value"):
                cards = [section]
            if cards:
                n = len(cards)
                card_w = min(55, 190 / n)
                gap = 3
                start_x = (210 - n * (card_w + gap) + gap) / 2
                y_start = pdf.get_y()
                for i, card in enumerate(cards):
                    x = start_x + i * (card_w + gap)
                    # Card background
                    pdf.set_fill_color(245, 245, 250)
                    pdf.rect(x, y_start, card_w, 40, "F")
                    # Accent top bar
                    pdf.set_fill_color(*acc_rgb)
                    pdf.rect(x, y_start, card_w, 3, "F")
                    # Value
                    pdf.set_xy(x, y_start + 6)
                    pdf.set_font(font_family, "B", 22)
                    pdf.set_text_color(*acc_rgb)
                    pdf.cell(card_w, 12, card.get("value", "0"), align="C")
                    # Label
                    pdf.set_xy(x, y_start + 20)
                    pdf.set_font(font_family, "", 9)
                    pdf.set_text_color(*muted_rgb)
                    pdf.cell(card_w, 8, card.get("label", ""), align="C")
                pdf.set_y(y_start + 48)

        elif sec_type == "callout":
            # Highlighted callout box with accent-tinted background
            callout_bg = _lighten(color_scheme.get("accent", "#1a1a2e"), 180)
            pdf.set_fill_color(*callout_bg)
            y_start = pdf.get_y()
            pdf.set_x(15)
            pdf.set_font(font_family, "", 11)
            pdf.set_text_color(*acc_rgb)
            pdf.multi_cell(180, 7, section.get("text", ""), fill=True)
            pdf.ln(5)

    # ── Save ──────────────────────────────────────────────────────────────
    pdf.output(str(output_path))
    page_count = pdf.pages_count
    return f"PDF saved: {output_path.name} ({page_count} pages)"


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN ORCHESTRATOR
# ═══════════════════════════════════════════════════════════════════════════════

_CREATION_PROMPT = """You are JARVIS's Document Design Engine. Given a user request, generate a JSON plan
for creating the requested document/visual. Return ONLY valid JSON, no markdown, no explanation.

OUTPUT FORMAT depends on the requested type:

--- PPTX ---
{
  "type": "pptx",
  "colors": {"bg": "#1a1a2e", "accent": "#e94560", "text": "#ffffff", "muted": "#aaaaaa"},
  "slides": [
    {"type": "title", "title": "MAIN TITLE", "subtitle": "Subtitle text"},
    {"type": "bullets", "title": "Section Title", "bullets": ["Point 1", "Point 2", "Point 3"]},
    {"type": "stats", "title": "Key Numbers", "stats": [{"value": "99%", "label": "Accuracy"}]},
    {"type": "two_column", "title": "Comparison", "left": ["Left items..."], "right": ["Right items..."]},
    {"type": "quote", "text": "Inspiring quote text", "author": "Author Name"},
    {"type": "section_divider", "title": "Section Name", "subtitle": "Optional subtitle", "section_number": "01"},
    {"type": "image", "title": "Visual Slide", "caption": "Description"}
  ]
}

--- XLSX ---
{
  "type": "xlsx",
  "colors": {"header_bg": "#1a1a2e", "header_text": "#ffffff", "accent": "#e94560", "alt_row": "#f0f0f5"},
  "sheets": [
    {
      "name": "Sheet Name",
      "headers": ["Col1", "Col2", "Col3"],
      "rows": [["val1", "val2", "val3"], ...],
      "chart": {"type": "bar|pie|line", "title": "Chart Title", "category_column": 1, "value_column": 2}
    }
  ]
}

--- DOCX ---
{
  "type": "docx",
  "colors": {"heading": "#1a1a2e", "accent": "#e94560", "body": "#333333"},
  "elements": [
    {"type": "heading", "text": "Title", "level": 1},
    {"type": "paragraph", "text": "Body text here..."},
    {"type": "bullet", "text": "Bullet point"},
    {"type": "numbered", "text": "Numbered item"},
    {"type": "table", "headers": ["Col1", "Col2"], "rows": [["a", "b"]]},
    {"type": "quote", "text": "Quote text"},
    {"type": "divider"},
    {"type": "image", "caption": "Image caption"}
  ]
}

--- POSTER ---
{
  "type": "poster",
  "width": 2480, "height": 3508,
  "colors": {"bg": "#1a1a2e", "accent": "#e94560", "text": "#ffffff", "muted": "#aaaaaa"},
  "elements": [
    {"type": "title", "text": "BIG TITLE", "x": 1240, "y": 200, "font_size": 120},
    {"type": "subtitle", "text": "Subtitle", "x": 1240, "y": 400, "font_size": 48},
    {"type": "text", "text": "Body text with word wrap", "x": 200, "y": 600, "font_size": 36, "max_width": 2000},
    {"type": "rectangle", "x": 100, "y": 1500, "width": 2280, "height": 400, "color": "#e94560", "radius": 20},
    {"type": "circle", "x": 500, "y": 800, "radius": 150, "color": "#00e5ff"},
    {"type": "line", "x1": 200, "y1": 1200, "x2": 2280, "y2": 1200, "color": "#e94560", "thickness": 4},
    {"type": "image", "x": 600, "y": 1800, "width": 1200, "height": 800}
  ]
}

--- INFOGRAPHIC ---
{
  "type": "infographic",
  "width": 1200, "height": 2400,
  "colors": {"bg": "#0f0f23", "accent": "#e94560", "secondary": "#00e5ff", "text": "#ffffff", "muted": "#888899"},
  "elements": [
    {"type": "title", "text": "INFOGRAPHIC TITLE", "x": 600, "y": 80, "font_size": 64},
    {"type": "stat_card", "x": 100, "y": 300, "width": 320, "height": 180, "value": "99%", "label": "Accuracy"},
    {"type": "divider", "y": 520},
    {"type": "text", "text": "Section text", "x": 100, "y": 550, "font_size": 28}
  ],
  "charts": [
    {
      "type": "bar|pie|line",
      "x": 100, "y": 700, "width": 500, "height": 350,
      "title": "Chart Title",
      "data": {"labels": ["A", "B", "C"], "values": [30, 50, 20]}
    }
  ]
}

--- LOGO ---
{
  "type": "logo",
  "size": 1024,
  "transparent_bg": false,
  "colors": {"bg": "#1a1a2e", "accent": "#e94560", "text": "#ffffff"},
  "elements": [
    {"type": "shape", "shape": "circle|rounded_rect|hexagon", "x": 512, "y": 400, "radius": 300, "color": "#e94560"},
    {"type": "ring", "x": 512, "y": 400, "radius": 350, "thickness": 10, "color": "#ffffff"},
    {"type": "initials", "text": "J", "x": 512, "y": 380, "font_size": 280, "color": "#ffffff"},
    {"type": "text", "text": "BRAND NAME", "x": 512, "y": 750, "font_size": 60, "color": "#ffffff"}
  ]
}

--- PDF ---
{
  "type": "pdf",
  "colors": {"bg": "#ffffff", "accent": "#1a1a2e", "text": "#333333", "muted": "#888888"},
  "date": "June 2026",
  "header": "Optional running header",
  "footer": "Optional footer (leave empty for page numbers)",
  "cover": {
    "title": "DOCUMENT TITLE",
    "subtitle": "Subtitle or description",
    "author": "Author Name",
    "stripe_color": "#e94560"
  },
  "sections": [
    {"type": "heading", "text": "Section Title", "level": 1},
    {"type": "text", "text": "Body paragraph text here..."},
    {"type": "bullet", "items": ["Bullet point 1", "Bullet point 2"]},
    {"type": "numbered", "items": ["Step 1", "Step 2"]},
    {"type": "table", "headers": ["Col1", "Col2"], "rows": [["a", "b"], ["c", "d"]]},
    {"type": "quote", "text": "Inspirational quote"},
    {"type": "stat_card", "cards": [{"value": "99%", "label": "Accuracy"}, {"value": "1.2K", "label": "Users"}]},
    {"type": "callout", "text": "Important highlighted note"},
    {"type": "divider"},
    {"type": "image", "caption": "Image description"}
  ]
}

RULES:
- Generate rich, professional, detailed content (not placeholders).
- For presentations: create 5-10 slides with varied layouts.
- For spreadsheets: generate realistic data rows (10-20 rows).
- For documents: create structured content with headings, paragraphs, bullets, tables.
- For posters: use large fonts, bold colors, geometric shapes.
- For infographics: include 2-3 charts with real data + stat cards.
- For logos: create a bold, clean design with shape + initials + brand name.
- Adapt the color scheme to match the topic (e.g., corporate=blues, food=reds/oranges).
- Use the user's image reference if provided.
"""


def doc_creator(
    parameters: dict = None,
    player=None,
    speak=None,
) -> str:
    """
    Jarvis Tool: Universal Document & Visual Creator.

    Creates presentations, spreadsheets, documents, posters, infographics, and logos.
    """
    params = parameters or {}
    prompt  = params.get("prompt", "").strip()
    doc_type = params.get("doc_type", "").strip().lower()
    image_path_raw = params.get("image_path", "").strip()

    if not prompt:
        return "Please provide a prompt describing what to create."

    # ── Validate image reference — only accept actual image files ──────
    image_ref = None
    if image_path_raw:
        try:
            _ip = Path(image_path_raw).resolve()
            if _ip.exists() and _ip.suffix.lower() in _IMAGE_EXTS:
                image_ref = _ip
        except Exception:
            pass

    # ── Extract content from non-image uploaded files (md, txt, csv, etc.) ──
    file_content = ""
    if image_path_raw:
        try:
            _fp = Path(image_path_raw).resolve()
            if _fp.exists() and _fp.suffix.lower() not in _IMAGE_EXTS:
                # Read text-based files
                if _fp.suffix.lower() in (".md", ".txt", ".csv", ".json", ".xml", ".html", ".py", ".js"):
                    file_content = _fp.read_text(encoding="utf-8", errors="replace")[:15000]
                elif _fp.suffix.lower() in (".pdf", ".docx", ".pptx", ".xlsx"):
                    # Try to extract text from Office/PDF files
                    file_content = _extract_file_text(_fp)
        except Exception:
            pass

    if player:
        player.write_log("[DocCreator] 🧬 Initializing Document Design Engine...")

    # Auto-detect type from prompt if not specified
    if not doc_type:
        prompt_lower = prompt.lower()
        if any(k in prompt_lower for k in ("presentation", "slide", "ppt", "pptx")):
            doc_type = "pptx"
        elif any(k in prompt_lower for k in ("spreadsheet", "sheet", "excel", "xlsx")):
            doc_type = "xlsx"
        elif any(k in prompt_lower for k in ("pdf",)):
            doc_type = "pdf"
        elif any(k in prompt_lower for k in ("document", "word", "doc", "docx", "report", "letter")):
            doc_type = "docx"
        elif any(k in prompt_lower for k in ("infographic", "data visualization", "chart")):
            doc_type = "infographic"
        elif any(k in prompt_lower for k in ("logo", "brand mark", "emblem")):
            doc_type = "logo"
        elif any(k in prompt_lower for k in ("poster", "flyer", "banner", "brochure")):
            doc_type = "poster"
        else:
            doc_type = "pptx"  # Default to presentation

    # Build prompt for Gemini
    image_context = ""
    if image_ref:
        image_context = f"\nThe user provided a reference image at: {image_ref}\nIncorporate this image into the design where appropriate."

    file_context = ""
    if file_content:
        file_context = f"\n\nThe user uploaded a file with the following content. Use this content as the basis for creating the document:\n\n--- FILE CONTENT START ---\n{file_content}\n--- FILE CONTENT END ---"

    full_prompt = f"{prompt}\n\nDocument type: {doc_type}{image_context}{file_context}\n\nGenerate the complete JSON design plan."

    if player:
        player.write_log(f"[DocCreator] 📐 Generating {doc_type.upper()} design plan...")

    # Get plan from Gemini (with retries & fallback models handled inside _gemini_generate)
    try:
        raw_response = _gemini_generate(
            prompt=full_prompt,
            system=_CREATION_PROMPT,
            model="gemini-2.5-flash"
        )
        clean_json = raw_response.strip()
        clean_json = re.sub(r"```(?:json)?", "", clean_json).strip().rstrip("`").strip()
        plan = json.loads(clean_json)
    except json.JSONDecodeError as e:
        snippet = raw_response[:500] if raw_response else "(empty)"
        return (
            f"Failed to parse design plan as JSON: {e}\n"
            f"Raw response snippet: {snippet}..."
        )
    except Exception as e:
        exc_name = type(e).__name__
        msg = str(e)
        if "503" in msg or "UNAVAILABLE" in msg.upper():
            return (
                f"Gemini API temporarily unavailable (503). "
                f"All fallback models also failed. Please try again in a minute.\n"
                f"Details: {exc_name}: {msg}"
            )
        elif "429" in msg or "RESOURCE_EXHAUSTED" in msg.upper():
            return (
                f"Gemini API rate limit hit (429). Please wait a moment and try again.\n"
                f"Details: {exc_name}: {msg}"
            )
        else:
            return f"Gemini planning failed ({exc_name}): {msg}"

    # Generate the actual file
    save_dir = _save_dir()
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    ext_map = {
        "pptx": ".pptx", "xlsx": ".xlsx", "docx": ".docx", "pdf": ".pdf",
        "poster": ".png", "infographic": ".png", "logo": ".png"
    }
    ext = ext_map.get(doc_type, ".pptx")
    output_path = save_dir / f"doc_creator_{timestamp}{ext}"

    if player:
        player.write_log(f"[DocCreator] 🎨 Rendering {doc_type.upper()}...")

    creators = {
        "pptx":        lambda: _create_pptx(plan, output_path, image_ref),
        "xlsx":        lambda: _create_xlsx(plan, output_path),
        "docx":        lambda: _create_docx(plan, output_path, image_ref),
        "pdf":         lambda: _create_pdf(plan, output_path, image_ref),
        "poster":      lambda: _create_poster(plan, output_path, image_ref),
        "infographic": lambda: _create_infographic(plan, output_path),
        "logo":        lambda: _create_logo(plan, output_path),
    }

    creator = creators.get(doc_type)
    if not creator:
        return f"Unknown document type: '{doc_type}'. Supported: pptx, xlsx, docx, pdf, poster, infographic, logo"

    try:
        result = creator()
    except Exception as e:
        import traceback
        traceback.print_exc()
        return f"Document creation failed: {e}"

    # Auto-open
    if output_path.exists():
        _open_file(output_path)

    return f"Sir, your {doc_type.upper()} has been created. {result}\nSaved to: Desktop/JarvisMedia/{output_path.name}"
